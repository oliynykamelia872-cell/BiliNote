import io
import json
import zipfile

import fitz
import pytest
from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation

from app.services.document_converter import DocumentConversionError, DocumentConverter


def _png_bytes() -> bytes:
    image = Image.new("RGB", (20, 20), "white")
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def test_docx_extracts_headings_tables_and_images(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    image_path = tmp_path / "picture.png"
    image_path.write_bytes(_png_bytes())
    document = Document()
    document.add_heading("项目标题", level=1)
    document.add_paragraph("正文内容")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "名称"
    table.cell(0, 1).text = "值"
    table.cell(1, 0).text = "状态"
    table.cell(1, 1).text = "完成"
    document.add_picture(str(image_path))
    source = tmp_path / "sample.docx"
    document.save(source)

    result = DocumentConverter("docx-task", source, ocr_mode="off").convert()

    assert "项目标题" in result.markdown
    assert "名称" in result.markdown
    assert result.assets
    assert list((tmp_path / "static" / "document_assets" / "docx-task").iterdir())


def test_pdf_extracts_text_and_images(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    source = tmp_path / "sample.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "PDF readable text")
    page.insert_image(fitz.Rect(72, 100, 110, 138), stream=_png_bytes())
    pdf.save(source)
    pdf.close()

    result = DocumentConverter("pdf-task", source, ocr_mode="off").convert()

    assert "PDF readable text" in result.markdown
    assert result.assets


def test_scanned_pdf_uses_offline_ocr(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    source = tmp_path / "scan.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_image(page.rect, stream=_png_bytes())
    pdf.save(source)
    pdf.close()
    converter = DocumentConverter("scan-task", source)
    monkeypatch.setattr(converter, "_offline_ocr", lambda _: "识别出的扫描文字")

    assert "识别出的扫描文字" in converter.convert().markdown


def test_rejects_unsupported_extension(tmp_path):
    source = tmp_path / "not-a-document.exe"
    source.write_text("hello", encoding="utf-8")

    with pytest.raises(DocumentConversionError, match="不支持"):
        DocumentConverter("bad-task", source).convert()


def test_markitdown_converts_presentation_spreadsheet_and_data(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[1]).shapes.title.text = "演示标题"
    pptx = tmp_path / "sample.pptx"
    presentation.save(pptx)
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.append(["名称", "数量"])
    worksheet.append(["苹果", 3])
    xlsx = tmp_path / "sample.xlsx"
    workbook.save(xlsx)
    data = tmp_path / "sample.json"
    data.write_text(json.dumps({"name": "BiliNote"}), encoding="utf-8")

    assert "演示标题" in DocumentConverter("ppt-task", pptx, ocr_mode="off").convert().markdown
    assert "苹果" in DocumentConverter("xlsx-task", xlsx, ocr_mode="off").convert().markdown
    assert "BiliNote" in DocumentConverter("json-task", data, ocr_mode="off").convert().markdown


def test_zip_converts_children_and_reports_failed_files(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    archive_path = tmp_path / "bundle.zip"
    with zipfile.ZipFile(archive_path, "w") as archive:
        archive.writestr("readme.txt", "压缩包正文")
        archive.writestr("unsupported.bin", b"binary")

    result = DocumentConverter("zip-task", archive_path, ocr_mode="off").convert()

    assert "压缩包正文" in result.markdown
    assert not result.failed_files


def test_rejects_unsafe_zip_paths(tmp_path):
    archive_path = tmp_path / "unsafe.zip"
    with zipfile.ZipFile(archive_path, "w") as archive:
        archive.writestr("../escape.txt", "unsafe")

    with pytest.raises(DocumentConversionError, match="不安全路径"):
        DocumentConverter("unsafe-task", archive_path, ocr_mode="off").convert()
